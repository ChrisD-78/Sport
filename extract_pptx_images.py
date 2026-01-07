#!/usr/bin/env python3
"""
Extrahiert Bilder aus einer PowerPoint-Datei
"""
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images_from_pptx(pptx_path, output_dir):
    """Extrahiert alle Bilder aus PowerPoint"""
    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)
    
    images = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape_num, shape in enumerate(slide.shapes, 1):
            # Prüfe ob es ein Bild ist
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                
                filename = f"slide_{slide_num}_img_{shape_num}.{image_ext}"
                filepath = os.path.join(output_dir, filename)
                
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)
                
                images.append({
                    'slide': slide_num,
                    'shape': shape_num,
                    'filename': filename,
                    'path': filepath
                })
                print(f"✓ Folie {slide_num}, Bild {shape_num}: {filename}")
            
            # Prüfe auch auf Gruppen (können Bilder enthalten)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = sub_shape.image
                        image_bytes = image.blob
                        image_ext = image.ext
                        
                        filename = f"slide_{slide_num}_group_{shape_num}_img_{len(images)+1}.{image_ext}"
                        filepath = os.path.join(output_dir, filename)
                        
                        with open(filepath, 'wb') as f:
                            f.write(image_bytes)
                        
                        images.append({
                            'slide': slide_num,
                            'shape': shape_num,
                            'filename': filename,
                            'path': filepath
                        })
                        print(f"✓ Folie {slide_num}, Gruppe {shape_num}, Bild: {filename}")
    
    return images

def main():
    pptx_path = "/Users/christofdrost/Desktop/2025_Landau LaOla_Skizze Förderantrag.pptx"
    output_dir = "/Users/christofdrost/Sitzungsvorlage/pptx_content/images"
    
    if not os.path.exists(pptx_path):
        print(f"Fehler: Datei nicht gefunden: {pptx_path}")
        return
    
    print(f"Extrahiere Bilder aus: {pptx_path}")
    images = extract_images_from_pptx(pptx_path, output_dir)
    
    if images:
        print(f"\n✓ {len(images)} Bilder extrahiert")
        print(f"Gespeichert in: {output_dir}")
    else:
        print("\n⚠ Keine Bilder gefunden")

if __name__ == "__main__":
    main()


