#!/usr/bin/env python3
"""
Konvertiert PowerPoint-Folien zu HTML-kompatiblen Bildern
"""
import os
from pptx import Presentation
from pptx.util import Inches
import base64

def extract_slide_content(pptx_path):
    """Extrahiert Inhalte aus PowerPoint und erstellt HTML"""
    prs = Presentation(pptx_path)
    slides_html = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_html = f'<div class="pptx-slide">\n'
        slide_html += f'<h3>Folie {i}</h3>\n'
        slide_html += '<div class="slide-content">\n'
        
        # Extrahiere Text-Inhalte
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_html += f'<p>{shape.text.strip()}</p>\n'
        
        slide_html += '</div>\n</div>\n'
        slides_html.append(slide_html)
    
    return slides_html

def main():
    pptx_path = "/Users/christofdrost/Desktop/2025_Landau LaOla_Skizze Förderantrag.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"Fehler: Datei nicht gefunden: {pptx_path}")
        return None
    
    slides_html = extract_slide_content(pptx_path)
    return slides_html

if __name__ == "__main__":
    result = main()
    if result:
        print(f"✓ {len(result)} Folien extrahiert")


