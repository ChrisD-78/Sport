#!/usr/bin/env python3
"""
Extrahiert alle Inhalte aus PowerPoint und erstellt HTML-Folien
"""
import os
from pptx import Presentation

def extract_all_content(pptx_path):
    """Extrahiert alle Inhalte aus PowerPoint"""
    prs = Presentation(pptx_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_data = {
            'number': i,
            'texts': [],
            'titles': []
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Versuche Titel zu identifizieren (meist größere/shape.text_frame.paragraphs[0].runs[0].font.size)
                if hasattr(shape, "text_frame"):
                    try:
                        if shape.text_frame.paragraphs:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                font_size = para.runs[0].font.size
                                if font_size and font_size > 20000000:  # Große Schrift = Titel
                                    slide_data['titles'].append(text)
                                else:
                                    slide_data['texts'].append(text)
                            else:
                                slide_data['texts'].append(text)
                        else:
                            slide_data['texts'].append(text)
                    except:
                        slide_data['texts'].append(text)
                else:
                    slide_data['texts'].append(text)
        
        slides_data.append(slide_data)
    
    return slides_data

def create_html_slides(slides_data):
    """Erstellt HTML-Code für die Folien"""
    html_slides = []
    
    for slide in slides_data:
        html = '<div class="slide pptx-integrated">\n'
        html += f'<h2>Förderantrag - Folie {slide["number"]}</h2>\n'
        html += '<div class="pptx-content">\n'
        
        if slide['titles']:
            for title in slide['titles']:
                html += f'<h3>{title}</h3>\n'
        
        if slide['texts']:
            html += '<ul class="pptx-text-list">\n'
            for text in slide['texts']:
                html += f'<li>{text}</li>\n'
            html += '</ul>\n'
        
        html += '</div>\n</div>\n'
        html_slides.append(html)
    
    return html_slides

def main():
    pptx_path = "/Users/christofdrost/Desktop/2025_Landau LaOla_Skizze Förderantrag.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"Fehler: Datei nicht gefunden: {pptx_path}")
        return None
    
    slides_data = extract_all_content(pptx_path)
    html_slides = create_html_slides(slides_data)
    
    # Speichere HTML in Datei
    output_file = "/Users/christofdrost/Sitzungsvorlage/pptx_slides.html"
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write('\n'.join(html_slides))
    
    print(f"✓ {len(html_slides)} Folien als HTML extrahiert")
    print(f"Gespeichert in: {output_file}")
    
    return html_slides

if __name__ == "__main__":
    main()


